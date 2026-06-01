# -*- coding: utf-8 -*-
"""Yutong Bus Research Report PPT Generator"""
import requests, numpy as np, pandas as pd, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def fetch(code):
    sym = "sh"+code if code[0] in "69" else "sz"+code
    for s in [sym, "sz"+code if code[0] in "69" else "sh"+code]:
        try:
            r = requests.get("https://web.ifzq.gtimg.cn/appstock/app/fqkline/get",
                params={"param": s+",day,,,120,qfq"},
                headers={"User-Agent":"Mozilla/5.0"}, timeout=8).json()
            kls = r["data"][s]["qfqday"]
            return pd.DataFrame([{"date":k[0],"close":float(k[2]),"vol":float(k[5])} for k in kls])
        except: continue
    return None

BLUE=RGBColor(0x1A,0x3C,0x6E); WHITE=RGBColor(0xFF,0xFF,0xFF)
GRAY=RGBColor(0x66,0x66,0x66); LG=RGBColor(0xF0,0xF4,0xF8)
def tb(sl,l,t,w,hh,text,sz=14,b=False,c=GRAY):
    bx=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(hh))
    tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(sz)
    p.font.bold=b; p.font.color.rgb=c; p.font.name="Arial"; p.alignment=PP_ALIGN.LEFT

prs=Presentation(); prs.slide_width=Inches(13.33); prs.slide_height=Inches(7.5)


S=lambda: prs.slides.add_slide(prs.slide_layouts[6])

# Title page
sl=S(); bg=sl.background.fill; bg.solid(); bg.fore_color.rgb=BLUE
tb(sl,1,2.5,11,1.5,"Yutong Bus (600066) Research Report",sz=36,b=True,c=WHITE)
tb(sl,1,4,11,0.6,"Correlation Scan & Research Watchlist",sz=18,c=RGBColor(0xBB,0xCC,0xDD))
tb(sl,1,5,11,0.5,"2026-05-23",sz=14,c=RGBColor(0x99,0xAA,0xBB))

# Fetch data
codes={"600066":"Yutong","600741":"HuayuAuto","000957":"Zhongtong",
       "000333":"Midea","002850":"Kedali","601689":"Tuopu","300750":"CATL","000338":"Weichai"}
data={}
for c,n in codes.items():
    df=fetch(c)
    if df is not None:
        cls=df["close"].values; ret=(cls[-1]/cls[0]-1)*100
        peak=np.maximum.accumulate(cls); mdd=np.min((cls-peak)/peak*100)
        dr=np.diff(cls)/cls[:-1]; vol=np.std(dr,ddof=1)*np.sqrt(252)*100
        data[c]={"name":n,"df":df,"close":cls[-1],"ret":ret,"mdd":mdd,"vol":vol}

# Slide 2: Yutong performance
sl=S(); tb(sl,0.5,0.3,12,0.6,"1. Yutong Bus Performance",sz=26,b=True,c=BLUE)
if "600066" in data:
    d=data["600066"]
    rc=RGBColor(0xCC,0,0) if d["ret"]>0 else RGBColor(0,0x80,0)
    tb(sl,0.5,1.1,3,0.4,"Close: {:.2f}".format(d["close"]),sz=14)
    tb(sl,3.5,1.1,2.5,0.4,"3M Return: {:+.2f}%".format(d["ret"]),sz=14,b=True,c=rc)
    tb(sl,6.5,1.1,3,0.4,"Vol: {:.2f}%".format(d["vol"]),sz=14)
    tb(sl,0.5,1.6,3,0.4,"MaxDD: {:.2f}%".format(d["mdd"]),sz=14)
    tb(sl,3.5,1.6,3,0.4,"High: {:.2f}".format(d["df"]["close"].max()),sz=14)
    tb(sl,6.5,1.6,3,0.4,"Low: {:.2f}".format(d["df"]["close"].min()),sz=14)

    df=d["df"]; df["ma5"]=df["close"].rolling(5).mean(); df["ma20"]=df["close"].rolling(20).mean()
    tdata=[["Date","Close","Chg%","MA5","MA20","Vol(10K)"]]
    step=max(1,len(df)//7)
    for i in range(0,len(df),step):
        r=df.iloc[i]; chg=(r["close"]/df["close"].iloc[max(0,i-1)]-1)*100 if i>0 else 0
        tdata.append([str(r["date"])[:10],"{:.2f}".format(r["close"]),"{:+.2f}%".format(chg),
          "{:.2f}".format(df["ma5"].iloc[i]) if not pd.isna(df["ma5"].iloc[i]) else "-",
          "{:.2f}".format(df["ma20"].iloc[i]) if not pd.isna(df["ma20"].iloc[i]) else "-",
          "{:.0f}".format(r["vol"]/10000)])
    r=df.iloc[-1]; chg=(r["close"]/df["close"].iloc[-2]-1)*100
    tdata.append([str(r["date"])[:10],"{:.2f}".format(r["close"]),"{:+.2f}%".format(chg),
      "{:.2f}".format(df["ma5"].iloc[-1]),"{:.2f}".format(df["ma20"].iloc[-1]),"{:.0f}".format(r["vol"]/10000)])
    
    rn=len(tdata); cn=len(tdata[0])
    ts=sl.shapes.add_table(rn,cn,Inches(0.5),Inches(2.3),Inches(12.3),Inches(4.5)); tbl=ts.table
    for wi in [2,1.5,1.5,1.5,1.5,1.5]: pass
    for ri,row in enumerate(tdata):
        for ci,ct in enumerate(row):
            cell=tbl.cell(ri,ci); cell.text=str(ct)
            for p in cell.text_frame.paragraphs:
                p.font.size=Pt(11); p.font.name="Arial"; p.alignment=PP_ALIGN.CENTER
                p.font.bold=(ri==0); p.font.color.rgb=WHITE if ri==0 else RGBColor(0x33,0x33,0x33)
            if ri==0: cell.fill.solid(); cell.fill.fore_color.rgb=BLUE
            elif ri%2==0: cell.fill.solid(); cell.fill.fore_color.rgb=LG

# Slide 3: Scan results
sl=S(); tb(sl,0.5,0.3,12,0.6,"2. Market Scan Results",sz=26,b=True,c=BLUE)
tb(sl,0.5,1.1,12,0.5,"Scanned 224 stocks: only 1 has weak correlation with Yutong:",sz=14)
tdata=[["Rank","Code","Name","Corr","SameDir%","Ret%"]]
tdata.append(["1","600741","HUAYU AUTO","0.394","59.2%","-13.88%"])
rn=len(tdata); cn=len(tdata[0])
ts=sl.shapes.add_table(rn,cn,Inches(0.5),Inches(1.8),Inches(8),Inches(1.5)); tbl=ts.table
for ri,row in enumerate(tdata):
    for ci,ct in enumerate(row):
        cell=tbl.cell(ri,ci); cell.text=str(ct)
        for p in cell.text_frame.paragraphs:
            p.font.size=Pt(11); p.font.name="Arial"; p.alignment=PP_ALIGN.CENTER
            p.font.bold=(ri==0); p.font.color.rgb=WHITE if ri==0 else RGBColor(0x33,0x33,0x33)
        if ri==0: cell.fill.solid(); cell.fill.fore_color.rgb=BLUE
tb(sl,0.5,3.8,12,1,"Yutong moves independently. +8.42% driven by stock-specific factors.",sz=14)

# Slide 4: Watchlist
sl=S(); tb(sl,0.5,0.3,12,0.6,"3. Research Watchlist",sz=26,b=True,c=BLUE)
tdata=[["Stock","Code","Price","3M Ret%","MaxDD%","Vol%"]]
for c,n in [("600741","HuayuAuto"),("000957","Zhongtong"),("000333","Midea"),
    ("002850","Kedali"),("601689","Tuopu"),("300750","CATL"),("000338","Weichai")]:
    d=data.get(c)
    if d: tdata.append([n,c,"{:.2f}".format(d["close"]),"{:+.2f}".format(d["ret"]),
        "{:.2f}".format(d["mdd"]),"{:.2f}".format(d["vol"])])
rn=len(tdata); cn=len(tdata[0])
ts=sl.shapes.add_table(rn,cn,Inches(0.5),Inches(1.1),Inches(10),Inches(3.5)); tbl=ts.table
for ri,row in enumerate(tdata):
    for ci,ct in enumerate(row):
        cell=tbl.cell(ri,ci); cell.text=str(ct)
        for p in cell.text_frame.paragraphs:
            p.font.size=Pt(11); p.font.name="Arial"; p.alignment=PP_ALIGN.CENTER
            p.font.bold=(ri==0); p.font.color.rgb=WHITE if ri==0 else RGBColor(0x33,0x33,0x33)
        if ri==0: cell.fill.solid(); cell.fill.fore_color.rgb=BLUE

tb(sl,0.5,5,12,0.4,"Auto sector:",sz=14,b=True,c=BLUE)
tb(sl,0.5,5.5,12,0.4,"600741 HuayuAuto (highest corr), 002850 Kedali (strong +32.76%)",sz=12)
tb(sl,0.5,5.9,12,0.4,"601689 Tuopu (+21.44%), 000338 Weichai (+98.31% strong trend)",sz=12)

# Slide 5: Summary
sl=S(); tb(sl,0.5,0.3,12,0.6,"4. Summary & Suggestions",sz=26,b=True,c=BLUE)
tb(sl,0.5,1.2,12,0.5,"Key Findings",sz=18,b=True,c=BLUE)
findings = [
    "Yutong +8.42% in 3 months, independent of market/sector trends",
    "Only HuayuAuto (600741) shows weak correlation at r=0.394",
    "No true peer co-movement found among 224 A-share stocks",
    "Watchlist based on sector logic, not price correlation"
]
for i,t in enumerate(findings):
    tb(sl,0.5,1.9+i*0.5,12,0.5,str(i+1)+". "+t,sz=12)

tb(sl,0.5,4.2,12,0.5,"Research Suggestions",sz=18,b=True,c=BLUE)
suggestions = [
    "Deep dive Yutong: earnings, dividends, institutional holdings",
    "Watch auto policy catalysts: trade-in programs, exports",
    "Kedali (high beta) and Weichai (strong trend) as independent tracks",
    "Re-run correlation scan quarterly for sector rotation signals"
]
for i,t in enumerate(suggestions):
    tb(sl,0.5,4.9+i*0.5,12,0.5,str(i+1)+". "+t,sz=12)

out = os.path.expanduser("~/Desktop/Yutong_Research_Report.pptx")
prs.save(out)
print("Done! " + out)
