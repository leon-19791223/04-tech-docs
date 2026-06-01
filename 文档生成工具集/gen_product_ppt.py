# -*- coding: utf-8 -*-
"""Product PPT: HCCDE-GaussDB Exam Practice System"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLUE = RGBColor(0x1A,0x3C,0x6E)
DARK = RGBColor(0x2D,0x2D,0x2D)
GRAY = RGBColor(0x66,0x66,0x66)
WHITE = RGBColor(0xFF,0xFF,0xFF)
ACCENT = RGBColor(0x00,0x7A,0xCC)
LGRAY = RGBColor(0xF0,0xF4,0xF8)
GREEN = RGBColor(0x10,0x7C,0x10)

def tb(sl,l,t,w,hh,text,sz=14,b=False,c=DARK,a=PP_ALIGN.LEFT,fn="Arial"):
    bx=sl.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(hh))
    tf=bx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(sz)
    p.font.bold=b; p.font.color.rgb=c; p.font.name=fn; p.alignment=a

def tbl(sl,l,t,w,hh,rows,cols=None):
    rn=len(rows); cn=len(rows[0])
    ts=sl.shapes.add_table(rn,cn,Inches(l),Inches(t),Inches(w),Inches(hh)); t=ts.table
    for ri,row in enumerate(rows):
        for ci,ct in enumerate(row):
            cell=t.cell(ri,ci); cell.text=str(ct)
            for p in cell.text_frame.paragraphs:
                p.font.size=Pt(10); p.font.name="Arial"; p.alignment=PP_ALIGN.CENTER
                p.font.bold=(ri==0); p.font.color.rgb=WHITE if ri==0 else DARK
            if ri==0: cell.fill.solid(); cell.fill.fore_color.rgb=BLUE
            elif ri%2==0: cell.fill.solid(); cell.fill.fore_color.rgb=LGRAY

S=lambda: prs.slides.add_slide(prs.slide_layouts[6])

# === Slide 1: Cover ===
sl=S(); bg=sl.background.fill; bg.solid(); bg.fore_color.rgb=BLUE
tb(sl,1,1.5,11,1.2,"HCCDE-GaussDB",sz=48,b=True,c=WHITE,a=PP_ALIGN.CENTER)
tb(sl,1,2.8,11,0.8,"Theory Exam Practice System",sz=28,c=RGBColor(0xBB,0xCC,0xDD),a=PP_ALIGN.CENTER)
tb(sl,1,3.8,11,0.6,"GaussDB 理论考题刷题系统",sz=20,c=RGBColor(0x99,0xAA,0xBB),a=PP_ALIGN.CENTER)
tb(sl,1,5,11,0.5,"Version 2.0 | 122 Questions | 6 Chapters",sz=14,c=RGBColor(0x88,0x99,0xAA),a=PP_ALIGN.CENTER)

# === Slide 2: Product Overview ===
sl=S()
tb(sl,0.5,0.3,12,0.6,"Product Overview",sz=28,b=True,c=BLUE)

# Feature cards as table
tdata=[["Feature","Description","Benefit"],
    ["Chapter-based Practice","Practice by chapter & question type","Targeted learning for weak areas"],
    ["Random Quiz","Random questions from full pool","Quick daily warm-up"],
    ["Mock Exam","50-question timed-style exam","Simulate real test conditions"],
    ["Auto Scoring","Instant correct/wrong feedback + score %","Track progress in real-time"],
    ["Error Review","Wrong answers saved with correct answers","Focused review after each session"],
    ["Reference Card","Key metrics quick reference","Last-minute revision aid"]]
tbl(sl,0.5,1.2,12.3,4.2,tdata)

# Stats box
tb(sl,0.5,5.8,3.5,0.5,"Question Bank: 122",sz=16,b=True,c=BLUE)
tb(sl,0.5,6.3,3.5,0.4,"Judge 30 | Single 60 | Multi 32",sz=12,c=GRAY)
tb(sl,4.5,5.8,3.5,0.5,"Chapters: 6",sz=16,b=True,c=BLUE)
tb(sl,4.5,6.3,3.5,0.4,"9%-28% weight distribution",sz=12,c=GRAY)
tb(sl,8.5,5.8,4,0.5,"Coverage: Full Exam Syllabus",sz=16,b=True,c=BLUE)
tb(sl,8.5,6.3,4,0.4,"Architecture, Security, Tuning, Ops",sz=12,c=GRAY)

# === Slide 3: Architecture ===
sl=S()
tb(sl,0.5,0.3,12,0.6,"System Architecture",sz=28,b=True,c=BLUE)

# Left: Data structure
tb(sl,0.5,1.2,6,0.5,"Data Organization",sz=18,b=True,c=DARK)
tdata=[["Layer","Component","Count"],
    ["Chapter","6 Chapters","Ch1-Ch6"],
    ["Question Type","Judge / Single / Multi","30+60+32"],
    ["Question Fields","ID / Q / Options / Ans / Analysis","5 fields each"],
    ["Pool","ALL_QUESTIONS (flat list)","122 items"]]
tbl(sl,0.5,1.8,6,2,tdata)

# Right: User flow
tb(sl,7,1.2,6,0.5,"User Flow",sz=18,b=True,c=DARK)
flow = [
    "1. Main Menu (5 options)",
    "   |",
    "2. Select Mode",
    "   |-- Chapter Practice -> Chapter -> Type",
    "   |-- Random Quiz -> Question count",
    "   |-- Mock Exam -> Auto 50 questions",
    "   |-- Stats / Reference -> View only",
    "   |",
    "3. Answer Loop",
    "   |-- Display question + options",
    "   |-- Accept input (√/×, A-D, combo)",
    "   |-- Show result + analysis",
    "   |-- Track correct/wrong",
    "   |",
    "4. Result Report",
    "   Score + Wrong answer review"
]
for i,line in enumerate(flow):
    c = GRAY if line.startswith(" ") else DARK
    b = "  " not in line
    tb(sl,7,1.7+i*0.3,6,0.3,line,sz=10,b=b,c=c)

# === Slide 4: Key Screens ===
sl=S()
tb(sl,0.5,0.3,12,0.6,"User Interface",sz=28,b=True,c=BLUE)

# Main menu mockup
tb(sl,0.5,1.1,6,0.5,"Main Menu",sz=16,b=True,c=BLUE)
menu_lines = [
    "========================================",
    "   HCCDE-GaussDB Quiz System v2.0",
    "   122 Questions | 6 Chapters",
    "========================================",
    "  [1] Chapter Practice",
    "  [2] Random Quiz",
    "  [3] Mock Exam",
    "  [4] Statistics",
    "  [5] Reference Card",
    "  [0] Exit",
    "========================================",
]
for i,line in enumerate(menu_lines):
    tb(sl,0.5,1.6+i*0.3,6,0.3,line,sz=10,c=DARK)

# Quiz screen mockup
tb(sl,7,1.1,6,0.5,"Quiz Screen",sz=16,b=True,c=BLUE)
quiz_lines = [
    "========================================",
    "  Q 3/5 [Ch3] [Single Choice]",
    "========================================",
    "  GaussDB default isolation level?",
    "",
    "  A. Read Uncommitted",
    "  B. Read Committed",
    "  C. Repeatable Read",
    "  D. Serializable",
    "",
    "  Your answer (A-D): _",
]
for i,line in enumerate(quiz_lines):
    tb(sl,7,1.6+i*0.3,6,0.3,line,sz=10,c=DARK)

# Result screen
tb(sl,0.5,5,6,0.5,"Result Report",sz=16,b=True,c=BLUE)
res_lines = [
    "========================================",
    "  Total: 10  Correct: 8  Wrong: 2",
    "  Score: 80.0%",
    "  Rating: **** Good!",
    "========================================",
    "  Wrong Answers:",
    "  [Ch3] D3-5: Your=A Correct=B",
    "  [Ch5] J5-2: Your=^ Correct=x",
]
for i,line in enumerate(res_lines):
    tb(sl,0.5,5.5+i*0.3,6,0.3,line,sz=10,c=DARK)

# Reference card
tb(sl,7,5,6,0.5,"Reference Card",sz=16,b=True,c=BLUE)
ref_lines = [
    "  Disaster Recovery Grades:",
    "  Grade  RTO        RPO",
    "  4     <=6h       <15min",
    "  5     min~2d     0~30min",
    "  6     seconds    0 (zero loss)",
    "",
    "  Deployment: Centralized vs Distributed",
    "  TPS <2万/Data <24TB -> Centralized",
]
for i,line in enumerate(ref_lines):
    tb(sl,7,5.5+i*0.3,6,0.3,line,sz=10,c=DARK)

# === Slide 5: Technical Highlights ===
sl=S()
tb(sl,0.5,0.3,12,0.6,"Technical Highlights",sz=28,b=True,c=BLUE)

highlights = [
    ("100% Python, Zero Dependencies",
     "Built entirely with Python standard library - no pip install required. Runs on any Python 3.6+ environment instantly."),
    ("Modular Question Bank Design",
     "Questions organized by chapter and type with consistent schema (id/question/options/answer/analysis). Easy to add new questions without touching code logic."),
    ("Three Learning Modes",
     "Chapter Practice (targeted), Random Quiz (daily drills), Mock Exam (simulation) - covering all learning stages from knowledge building to exam readiness."),
    ("Smart Scoring & Error Tracking",
     "Real-time correct/wrong feedback with detailed analysis. Session-level error tracking highlights knowledge gaps for focused review."),
    ("Console UI with Zero Config",
     "Clean terminal interface with box-drawing UI. No GUI framework, no web server, no database - just run and study."),
]

for i,(title,desc) in enumerate(highlights):
    tb(sl,0.5,1.2+i*1.15,12,0.4,str(i+1)+". "+title,sz=16,b=True,c=BLUE)
    tb(sl,0.5,1.6+i*1.15,12,0.5,desc,sz=12,c=GRAY)

# === Slide 6: Product Roadmap ===
sl=S()
tb(sl,0.5,0.3,12,0.6,"Product Roadmap",sz=28,b=True,c=BLUE)

tdata=[["Version","Timeline","Features"],
    ["v1.0","Current","6 chapters, 122 questions, 3 modes, console UI"],
    ["v2.0","Current+","Multi-session wrong book, progress persistence"],
    ["v2.5","Next","Question bank editor GUI, custom quiz builder"],
    ["v3.0","Future","Web UI, leaderboard, timed mock exams"],
    ["v3.5","Future","Mobile-friendly PWA, cloud sync"]]
tbl(sl,0.5,1.2,12.3,3,tdata)

tb(sl,0.5,4.8,12,0.5,"Current Focus Areas",sz=18,b=True,c=BLUE)
focus = [
    "Wrong book persistence (JSON file) - track mistakes across sessions",
    "Progress statistics by chapter - identify weak areas visually",
    "Timed exam mode - add countdown timer for real exam simulation",
    "Custom quiz builder - select specific chapters/types/question count",
]
for i,f in enumerate(focus):
    tb(sl,0.5,5.5+i*0.4,12,0.4,str(i+1)+". "+f,sz=12,c=GRAY)

# Save
out = os.path.expanduser("~/Desktop/HCCDE_GaussDB_Product_Showcase.pptx")
prs.save(out)
print("Saved: " + out)
