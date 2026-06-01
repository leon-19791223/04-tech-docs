# -*- coding: utf-8 -*-
"""补充数据湖仓一体化方案的 4.数据校验 和 5.Hudi小文件 内容"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import os, copy

# ── 颜色方案（匹配原PPT）──
BLUE_DARK = RGBColor(0x1A, 0x3C, 0x6E)
BLUE = RGBColor(0x00, 0x7A, 0xCC)
BLUE_LIGHT = RGBColor(0xE8, 0xF0, 0xF8)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x10, 0x7C, 0x10)
RED = RGBColor(0xCC, 0x33, 0x33)
ORANGE = RGBColor(0xE8, 0xA8, 0x20)
TABLE_HDR = BLUE_DARK
TABLE_ALT = RGBColor(0xF0, 0xF4, 0xF8)

path = os.path.expanduser('~/Desktop/数据湖仓一体化-实时同步解决方案-2026.pptx')
prs = Presentation(path)
SW = prs.slide_width   # 13.33"
SH = prs.slide_height  # 7.5"

# ── Helper functions ──
def add_textbox(slide, left, top, width, height, text, sz=14, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, font_name="微软雅黑"):
    bx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font_name
    p.alignment = align
    return tf

def add_multiline(slide, left, top, width, height, lines, sz=10, color=DARK):
    """Add multiple lines of text in a single textbox"""
    bx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = bx.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Determine if it's a header line (starts with special markers)
        is_header = line.startswith("■") or line.startswith("◆") or line.startswith("┌") or line.startswith("├") or line.startswith("└")
        p.text = line
        p.font.size = Pt(sz)
        p.font.bold = is_header
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.space_after = Pt(2)
    return tf

def add_table(slide, left, top, width, height, rows, col_widths=None):
    """Add a styled table matching original PPT style"""
    rn = len(rows); cn = len(rows[0])
    ts = slide.shapes.add_table(rn, cn, Inches(left), Inches(top), Inches(width), Inches(height))
    tbl = ts.table
    if col_widths:
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = Inches(w)

    for ri, row in enumerate(rows):
        for ci, ct in enumerate(row):
            cell = tbl.cell(ri, ci); cell.text = str(ct)
            # Cell padding
            tcPr = cell._tc.get_or_add_tcPr()
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)

            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9); p.font.name = "微软雅黑"
                p.alignment = PP_ALIGN.LEFT if ci > 0 else PP_ALIGN.CENTER
                p.font.bold = (ri == 0)
                p.font.color.rgb = WHITE if ri == 0 else DARK
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TABLE_HDR
            elif ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TABLE_ALT
    return ts

def add_title_bar(slide, text):
    """Add a title bar at top of slide"""
    # Blue background bar
    shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SW, Inches(0.75)
    )
    shape.fill.solid(); shape.fill.fore_color.rgb = BLUE_DARK
    shape.line.fill.background()
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(24); p.font.bold = True
    p.font.color.rgb = WHITE; p.font.name = "微软雅黑"
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.12)

def section_title(slide, left, top, text, sz=18):
    """Section sub-title"""
    add_textbox(slide, left, top, 12, 0.5, text, sz=sz, bold=True, color=BLUE_DARK)

def bullet_list(slide, left, top, items, sz=10, color=DARK, start_idx=None):
    """List of bullet items"""
    for i, item in enumerate(items):
        prefix = f"{start_idx+i}. " if start_idx else "  * "
        add_textbox(slide, left, top + i * 0.32, 12, 0.35,
                    prefix + item, sz=sz, color=color)

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[0])  # blank

# ════════════════════════════════════════════════════════════════
# SLIDE 7: 数据校验与对账 — 概述
# ════════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, "补充点4：数据校验与对账机制")

# 左侧：为什么需要
section_title(sl, 0.5, 1.0, "为什么需要数据校验")
add_multiline(sl, 0.5, 1.5, 6.0, 1.2, [
    "CDC实时同步链路中，数据可能在任一环节丢失或错乱：",
    "  源库 → Z-RPS(CDC) → Kafka → Octopus(Flink) → Hudi/Hive",
    "金融场景下，一条数据不一致就可能导致监管报送异常。",
], sz=11)

# 分层校验架构图（右）
section_title(sl, 7.0, 1.0, "三层校验架构")
tier_data = [
    ["层级", "校验方式", "频次", "覆盖范围"],
    ["实时校验层", "逐行CRC / 主键去重 / 延迟监控", "秒级", "Flink 内嵌算子"],
    ["定时校验层", "分区行数比对 / 关键字段SUM", "小时级", "Kafka侧 vs 目标侧"],
    ["全量校验层", "全表字段级比对 / 自动修复", "天级", "T+1 凌晨窗口"],
]
add_table(sl, 7.0, 1.5, 6.0, 1.8, tier_data)

# 校验技术方案
section_title(sl, 0.5, 3.5, "技术方案")
tech_items = [
    "实时校验：Flink 作业内嵌 CRC 算子，每条数据计算 MD5(data_crc) 随同步流写入",
    "对账查询：源 vs 目标 FULL JOIN 按主键比对 data_crc，差异即告警",
    "定时校验：每小时对每个分区 COUNT(*) + SUM(关键金额字段)，阈值 0.1%",
    "全量校验：DataX 分片并发（按主键 hash 16-32片），逐字段比差异",
    "修复流程：差异数据自动生成修复SQL -> 人工审批 -> 执行重跑",
]
bullet_list(sl, 0.5, 4.0, tech_items, sz=10)

# 底部：关键指标
section_title(sl, 0.5, 5.8, "金融场景关键 SLA 指标")
sla_data = [
    ["指标", "建议值", "说明"],
    ["RPO（数据丢失量）", "≤ 5秒", "超过自动告警"],
    ["数据准确率", "99.999%", "月度全量校验通过率"],
    ["对账覆盖度", "100%", "所有同步表必须配置校验"],
    ["修复时效", "< 30分钟", "从发现到修复完成"],
]
add_table(sl, 0.5, 6.2, 8.0, 1.4, sla_data)

# ════════════════════════════════════════════════════════════════
# SLIDE 8: 对账架构图
# ════════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, "数据校验与对账 — 架构集成")

# Architecture flow
add_multiline(sl, 0.5, 1.0, 12, 3.5, [
    "■ 对账系统架构",
    "",
    "  源端 Oracle/MySQL ──→ Z-RPS ──→ Kafka ──→ Octopus(Flink) ──→ 目标 DWS/MRS",
    "                            ↓                   ↓",
    "                    Kafka侧校验Topic        Flink侧输出校验结果",
    "                            ↓                   ↓",
    "                    ┌──────────────────────────────────────────┐",
    "                    │          对账告警中心                      │",
    "                    │  Prometheus + AlertManager               │",
    "                    │  企业微信 / 钉钉 / 短信告警               │",
    "                    └──────────────────────────────────────────┘",
    "",
    "◆ 实时校验：Flink SQL 追加 MD5(data_crc) + CURRENT_TIMESTAMP",
    "◆ 定时校验：每小时聚合 COUNT(*), SUM(amount)，差异>0.1%告警",
    "◆ 全量校验：T+1 凌晨窗口 DataX 分片并行比对",
    "◆ 告警方式：延迟>30s / 行数不匹配 / CRC不一致 → 即时通知",
], sz=11, color=DARK)

# 推荐工具
section_title(sl, 0.5, 4.8, "推荐对账工具矩阵")
tools_data = [
    ["场景", "推荐工具", "部署方式", "备注"],
    ["实时行数校验", "Flink + Kafka", "随同步作业部署", "开销低，秒级延迟"],
    ["定时字段校验", "自研 Python 脚本", "独立定时任务", "灵活定制校验规则"],
    ["全量数据对账", "DataX / 华为 CDM", "独立调度", "分片并发，性能好"],
    ["一致性校验", "华为 DataArts", "云原生", "对接 MRS/DWS 原生"],
    ["告警通知", "AlertManager", "独立部署", "支持多渠道推送"],
]
add_table(sl, 0.5, 5.2, 8.5, 2.0, tools_data)

# ════════════════════════════════════════════════════════════════
# SLIDE 9: Hudi 小文件问题 — 概述
# ════════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, "补充点5：Flink实时写入Hudi的小文件问题")

# 问题描述
section_title(sl, 0.5, 1.0, "问题本质")
add_multiline(sl, 0.5, 1.4, 12, 1.5, [
    "Flink 流处理以 checkpoint 间隔为批次写入 Hudi，每个批次每分区至少生成一个文件：",
    "  假设：checkpoint = 60s，分区有10个bucket，写入8小时",
    "  → 文件数 = 8 × 60 ÷ 1 × 10 = 4800 个 parquet 文件",
    "  → 每个文件仅几MB（远小于HDFS默认块大小128MB）",
], sz=11)

# 危害
section_title(sl, 0.5, 2.6, "小文件的危害")
harm_data = [
    ["危害", "影响", "后果"],
    ["NameNode 元数据压力", "每个文件一条元数据", "NN 内存暴涨，GC 频繁"],
    ["查询性能下降", "扫描需打开数千个小文件", "查询延迟增加 5-10 倍"],
    ["存储利用率低", "文件 << HDFS 块大小(128MB)", "浪费 90%+ 存储空间"],
    ["Task 调度开销", "每个文件一个 Split/InputSplit", "Spark/Flink 调度瓶颈"],
]
add_table(sl, 0.5, 3.0, 10.0, 1.6, harm_data)

# 解决方案概览
section_title(sl, 0.5, 5.0, "解决方案四步法")
solutions = [
    "方案一：合理配置 Compaction + Clustering 参数（MOR表合并 + 小文件合并）",
    "方案二：Flink 参数调优（增大 checkpoint 间隔、batch size）",
    "方案三：分阶段合并策略（白天写入 → 凌晨 Clustering 合并为 128MB）",
    "方案四：优化分区策略（避免分钟级分区，使用 Bucket 索引）",
]
bullet_list(sl, 0.5, 5.5, solutions, sz=11, start_idx=1)

# ════════════════════════════════════════════════════════════════
# SLIDE 10: Hudi 参数配置详解
# ════════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, "Hudi 小文件 — 参数配置详解")

# 左侧：Hudi 表参数
section_title(sl, 0.5, 1.0, "Hudi 表属性配置")
hudi_params = [
    ["参数", "推荐值", "说明"],
    ["hoodie.compact.inline", "true", "启用在线 Compaction（MOR表）"],
    ["hoodie.compact.inline.max.delta.commits", "5", "每5个commit触发一次Compaction"],
    ["hoodie.clustering.inline", "true", "启用在线 Clustering"],
    ["hoodie.clustering.inline.max.commits", "4", "每4个commit触发"],
    ["hoodie.parquet.max.file.size", "134217728", "目标文件大小 128MB"],
    ["hoodie.parquet.small.file.limit", "134217728", "小于此值触发合并"],
    ["hoodie.clustering.plan.strategy.target.file.max.bytes", "134217728", "合并目标大小"],
    ["hoodie.index.type", "BUCKET", "Bucket索引替代Bloom索引"],
    ["hoodie.bucket.index.num.buckets", "64", "Bucket数量"],
]
add_table(sl, 0.5, 1.5, 6.5, 2.8, hudi_params)

# 右侧：Flink 参数
section_title(sl, 7.5, 1.0, "Flink 写入参数")
flink_params = [
    ["参数", "推荐值", "说明"],
    ["checkpoint.interval", "120s", "增大间隔，减少写入批次"],
    ["checkpoint.min-pause", "60s", "最小暂停间隔"],
    ["write.tasks", "4", "写入并行度（不要过高）"],
    ["write.batch.size", "512 MB", "批次大小"],
    ["write.operation", "upsert", "upsert 替代 insert"],
    ["table.exec.sink.upsert-materialize", "NONE", "减少数据膨胀"],
    ["pipeline.operator-chaining", "true", "算子链优化"],
]
add_table(sl, 7.5, 1.5, 5.5, 2.2, flink_params)

# 分阶段策略
section_title(sl, 0.5, 4.6, "分阶段合并策略")

add_multiline(sl, 0.5, 5.0, 12, 2.5, [
    "┌──────────────────────┐       ┌──────────────────────────┐",
    "│  写入阶段 (7×24)      │       │  合并阶段 (凌晨低峰)      │",
    "│  Flink 实时写入        │  自动  │  Clustering 合并         │",
    "│  小文件 ~16-32MB      │──触发─→│  目标文件 → 128MB        │",
    "│  (快速写入，低延迟)    │  异步  │  (查询优化，存储高效)     │",
    "└──────────────────────┘       └──────────────────────────┘",
    "",
    "◆ 白天：hoodie.parquet.max.file.size = 33554432 (32MB) — 写得快",
    "◆ 凌晨：触发 Clustering 合并为 128MB — 查得快",
    "◆ 监控指标：单分区文件数 > 500 告警，平均文件大小 < 32MB 告警",
], sz=10, color=DARK)

# ════════════════════════════════════════════════════════════════
# SLIDE 11: 分区策略对比 + 总结
# ════════════════════════════════════════════════════════════════
sl = new_slide()
add_title_bar(sl, "分区策略优化与总结")

# 分区策略对比
section_title(sl, 0.5, 1.0, "分区策略优化（从源头控制小文件）")
part_data = [
    ["策略", "示例", "文件数/天", "写入性能", "查询性能", "推荐度"],
    ["分钟级分区 (❌)", "dt/yyyy/MM/dd/HH/mm", "> 10000", "差", "差", "强烈不推荐"],
    ["小时级分区", "dt/yyyy/MM/dd/HH", "~720", "中", "中", "不推荐"],
    ["天级分区", "dt=20260528", "~12", "好", "好", "推荐"],
    ["天+Bucket(64)", "dt=... + bucket索引", "~64", "好", "很好", "最推荐"],
]
add_table(sl, 0.5, 1.5, 12, 1.8, part_data)

# 建议配置总结
section_title(sl, 0.5, 3.8, "推荐配置清单")

summary_data = [
    ["配置项", "推荐配置", "效果"],
    ["Checkpoint 间隔", "120s（不低于60s）", "减少文件生成频率"],
    ["分区策略", "dt/hh（天+小时）", "避免分钟级分区"],
    ["Clustering", "每5个commit触发一次", "自动合并小文件"],
    ["目标文件大小", "128MB", "与HDFS块大小匹配"],
    ["写入模式", "upsert + Bucket索引", "减少数据膨胀"],
    ["凌晨合并窗口", "异步 Clustering", "不影响白天写入延迟"],
    ["并行度", "write.tasks=4", "过多并发加重小文件"],
]
add_table(sl, 0.5, 4.2, 10.0, 2.3, summary_data)

# 底部总结
add_textbox(sl, 0.5, 6.7, 12, 0.5,
    "总结：合理的 Flink Checkpoint + Hudi Clustering + 分区策略 可有效控制小文件数量，"
    "在写入性能和查询性能之间取得平衡",
    sz=12, bold=True, color=BLUE_DARK)

# ════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════
out = os.path.expanduser("~/Desktop/数据湖仓一体化-实时同步解决方案-2026.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
