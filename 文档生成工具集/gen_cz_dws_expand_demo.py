# -*- coding: utf-8 -*-
"""Generate Chouzhou Bank DWS Expansion Demo PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

BLUE_DARK = RGBColor(0x1A, 0x3C, 0x6E)
BLUE = RGBColor(0x00, 0x7A, 0xCC)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x10, 0x7C, 0x10)
RED = RGBColor(0xCC, 0x33, 0x33)
ORANGE = RGBColor(0xE8, 0xA8, 0x20)
TABLE_HDR = BLUE_DARK
TABLE_ALT = RGBColor(0xF0, 0xF4, 0xF8)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

S = lambda: prs.slides.add_slide(prs.slide_layouts[6])

def add_tb(sl, left, top, width, height, text, sz=14, bold=False,
           color=DARK, align=PP_ALIGN.LEFT, font_name="Arial"):
    bx = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font_name
    p.alignment = align
    return tf

def add_ml(sl, left, top, width, height, lines, sz=11, color=DARK):
    """Add multiple lines of text in one textbox"""
    bx = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = bx.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        txt = line.get("t", line) if isinstance(line, dict) else line
        fs = line.get("s", sz) if isinstance(line, dict) else sz
        fc = line.get("c", color) if isinstance(line, dict) else color
        b = line.get("b", False) if isinstance(line, dict) else False
        p.text = txt; p.font.size = Pt(fs); p.font.bold = b
        p.font.color.rgb = fc; p.font.name = "Arial"
        p.space_after = Pt(3)
    return tf

def add_rect(sl, left, top, width, height, fill_color):
    shp = sl.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp

def add_tbl(sl, left, top, width, height, rows, col_w=None):
    rn = len(rows); cn = len(rows[0])
    ts = sl.shapes.add_table(rn, cn, Inches(left), Inches(top), Inches(width), Inches(height))
    t = ts.table
    if col_w:
        for i, w in enumerate(col_w):
            t.columns[i].width = Inches(w)
    for ri, row in enumerate(rows):
        for ci, ct in enumerate(row):
            cell = t.cell(ri, ci); cell.text = str(ct)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9); p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER
                p.font.bold = (ri == 0)
                p.font.color.rgb = WHITE if ri == 0 else DARK
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TABLE_HDR
            elif ri % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TABLE_ALT
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
    return t

def cover_slide(title, subtitle, info):
    sl = S()
    bg = sl.background.fill; bg.solid(); bg.fore_color.rgb = BLUE_DARK
    add_tb(sl, 1, 1.5, 11, 1.2, title, sz=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tb(sl, 1, 2.8, 11, 0.8, subtitle, sz=28, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.CENTER)
    add_tb(sl, 1, 3.8, 11, 0.6, info, sz=16, color=RGBColor(0x99, 0xAA, 0xBB), align=PP_ALIGN.CENTER)
    add_tb(sl, 1, 5.0, 11, 0.5, "掌数科技 | 2026年5月", sz=14, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)
    return sl

def section_header(title):
    sl = S()
    add_tb(sl, 0.5, 0.3, 12, 0.6, title, sz=28, bold=True, color=BLUE_DARK)
    return sl

# ==================== SLIDES ====================

# S1: Cover
cover_slide(
    "稠州银行",
    "GaussDB(DWS) 集群扩容方案 Demo",
    "GaussDB(DWS) 8.1.3 x86 -> 8.2.1 aarch64  平滑扩容  |  数据重分布  |  风险控制"
)

# S2: Content Overview
sl = S()
add_tb(sl, 0.5, 0.3, 12, 0.6, "演示内容概览", sz=28, bold=True, color=BLUE_DARK)
items = [
    ("01", "项目背景与现状分析", "当前架构、瓶颈分析、扩容目标"),
    ("02", "扩容方案总体设计", "技术选型、架构对比、方案优势"),
    ("03", "扩容实施详细流程", "LLD规划、preinstall、添加节点、实例部署"),
    ("04", "数据重分布策略", "在线/离线重分布、并发控制、进度监控"),
    ("05", "风险控制与应急预案", "前置检查、巡检项、故障处理矩阵"),
    ("06", "回滚方案", "各阶段回退策略与操作步骤"),
    ("07", "项目实施计划", "阶段划分、时间线、交付物清单"),
]
for i, (num, title, desc) in enumerate(items):
    y = 1.2 + i * 0.85
    add_rect(sl, 0.5, y, 0.5, 0.5, BLUE_DARK)
    add_tb(sl, 0.5, y + 0.05, 0.5, 0.4, num, sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tb(sl, 1.3, y, 4, 0.35, title, sz=16, bold=True, color=DARK)
    add_tb(sl, 1.3, y + 0.35, 8, 0.3, desc, sz=11, color=GRAY)

# S3: Project Background
sl = S()
add_tb(sl, 0.5, 0.3, 12, 0.6, "1. 项目背景与现状分析", sz=26, bold=True, color=BLUE_DARK)
add_tb(sl, 0.5, 1.0, 12, 0.4, "◎当前架构概况", sz=16, bold=True, color=BLUE)
add_ml(sl, 0.5, 1.5, 12, 2.0, [
    {"t": "• 数仓系统底层数据库：GaussDB(DWS) 8.1.3 x86 版本", "s": 12},
    {"t": "• 集群拓扑：CMS主112 / 备119，多节点部署架构", "s": 12},
    {"t": "• 当前可堆积空间：约245TB，面临存储与性能瓶颈", "s": 12},
    {"t": "• 信创战略要求：需提升自主可控能力，向 aarch64 架构迁移", "s": 12},
])
add_tb(sl, 0.5, 3.5, 12, 0.4, "◎扩容驱动力", sz=16, bold=True, color=BLUE)
add_tbl(sl, 0.5, 4.0, 12.3, 2.8, [
    ["维度", "当前状态", "目标状态", "驱动因素"],
    ["硬件架构", "x86 架构", "aarch64 (鲲鹏)", "信创战略、自主可控"],
    ["数据库版本", "GaussDB(DWS) 8.1.3", "GaussDB(DWS) 8.2.1", "功能增强、性能优化、安全合规"],
    ["集群容量", "~245TB 可用", "按需扩容至500TB+", "业务数据年增长 ~30%"],
    ["高可用能力", "基础主备", "增强容灾 + Roach备份", "RPO/RTO 优化"],
    ["性能需求", "部分场景瓶颈", "内核优化 + 资源扩展", "复杂查询、高并发场景"],
], col_w=[2.5, 3, 3, 3.8])

# S4: Overall Architecture
sl = S()
add_tb(sl, 0.5, 0.3, 12, 0.6, "2. 扩容方案总体设计（双集群容灾迁移+扩容）", sz=24, bold=True, color=BLUE_DARK)
add_tb(sl, 0.5, 1.0, 12, 0.4, "◎方案架构：三步走策略 - 迁移 -> 升级 -> 扩容", sz=16, bold=True, color=BLUE)
# Phase boxes
phases = [
    ("Phase 1", "双集群迁移", "8.1.3 x86 -> 8.2.1 aarch64\nRoach工具, 全量+增量同步", BLUE_DARK),
    ("Phase 2", "目标集群升级", "8.2.1 补丁升级 & 服务重启\nUpdateService + 巡检", BLUE),
    ("Phase 3", "集群扩容", "LLD规划 + 添加节点 + 重分布\n在线/离线 模式", GREEN),
]
for i, (phase, title, desc, color) in enumerate(phases):
    x = 0.5 + i * 4.2
    add_rect(sl, x, 1.5, 3.8, 1.2, color)
    add_tb(sl, x + 0.2, 1.55, 3.4, 0.35, phase, sz=11, bold=True, color=WHITE)
    add_tb(sl, x + 0.2, 1.85, 3.4, 0.4, title, sz=16, bold=True, color=WHITE)
    add_tb(sl, x + 0.2, 2.2, 3.4, 0.4, desc, sz=10, color=RGBColor(0xDD, 0xEE, 0xFF))
    if i < 2:
        add_tb(sl, x + 3.8, 1.8, 0.4, 0.4, "->", sz=24, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

add_tb(sl, 0.5, 3.0, 12, 0.4, "◎核心技术选型", sz=16, bold=True, color=BLUE)
add_tbl(sl, 0.5, 3.5, 12.3, 3.2, [
    ["技术领域", "选型方案", "说明"],
    ["数据迁移工具", "Roach + SyncDataToStby.py", "GaussDB原生备份恢复工具，支持全量+增量同步"],
    ["集群容灾", "双集群容灾架构", "主集群提供读写，备集群只读热备，支持角色切换"],
    ["扩容工具", "LLD配置规划工具", "自动生成扩容模板（expandTemplet.xml）"],
    ["数据重分布", "在线/离线重分布", "支持在线模式（业务不中断），分段执行"],
    ["预检查工具", "SysChecker / FusionCare", "扩容前全面巡检，识别潜在风险"],
    ["升级工具", "UpdateService", "Manager化升级管理，支持灰度+回滚"],
], col_w=[2.5, 3.5, 6.3])

# S5: Topology
sl = S()
add_tb(sl, 0.5, 0.3, 12, 0.6, "2.2 扩容架构拓扑", sz=26, bold=True, color=BLUE_DARK)
add_tb(sl, 0.5, 1.0, 12, 0.4, "◎节点扩容前后对比", sz=16, bold=True, color=BLUE)
add_tbl(sl, 0.5, 1.5, 12.3, 2.5, [
    ["角色", "扩容前节点数", "扩容后节点数", "新增说明"],
    ["CMS (主/备)", "2", "2", "保持不变"],
    ["GTM (主/备)", "2", "2", "保持不变"],
    ["CN (Coordinator)", "≤10", "≤10", "需确保总数不超过10"],
    ["DN (Datanode)", "N", "N + M", "新增节点与原节点DN数一致"],
    ["管理节点(OMS)", "2", "2", "保持不变"],
], col_w=[3, 2.5, 2.5, 4.3])

add_tb(sl, 0.5, 4.3, 12, 0.4, "◎扩容数据流示意", sz=16, bold=True, color=BLUE)
flow_steps = [
    ("① LLD规划", "生成扩容模板\n配置IP/实例拓扑", BLUE_DARK),
    ("② Preinstall", "新节点OS配置\n磁盘分区/格式化", BLUE),
    ("③ 添加主机", "Manager导入模板\n安装NodeAgent", GREEN),
    ("④ 添加实例", "部署DN/CN实例\n配置同步", ORANGE),
    ("⑤ 数据重分布", "在线/离线模式\n数据均匀分布", RED),
]
for i, (step, desc, color) in enumerate(flow_steps):
    x = 0.5 + i * 2.5
    add_rect(sl, x, 4.8, 2.1, 1.8, color)
    add_tb(sl, x + 0.1, 4.9, 1.9, 0.5, step, sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tb(sl, x + 0.1, 5.4, 1.9, 1.0, desc, sz=10, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        add_tb(sl, x + 2.1, 5.4, 0.4, 0.4, "->", sz=20, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# S6: Implementation Detail
sl = S()
add_tb(sl, 0.5, 0.3, 12, 0.6, "3. 扩容实施详细流程", sz=26, bold=True, color=BLUE_DARK)
add_tb(sl, 0.5, 1.0, 12, 0.4, "◎前置条件检查", sz=16, bold=True, color=BLUE)
add_ml(sl, 0.5, 1.5, 12, 1.8, [
    {"t": "□ comm_max_datanode >= 扩容后DN总数（例：128->256） -> gs_guc set + 重启集群", "s": 11},
    {"t": "□ 执行 ANALYZE VERBOSE 保证统计信息准确性", "s": 11},
    {"t": "□ FusionInsight Manager中 MPPDB/KrbServer/LdapServer 运行状态为【良好】", "s": 11},
    {"t": "□ 检查并清除MPPDB相关告警", "s": 11},
    {"t": "□ 确认新节点OS编码格式为 en_US.UTF-8", "s": 11},
    {"t": "□ 确认镜像文件已挂载，数据盘已分区且无业务数据", "s": 11},
])
add_tb(sl, 0.5, 3.4, 12, 0.4, "◎LLD配置规划工具操作流程", sz=16, bold=True, color=BLUE)
lld_steps = [
    ("1", "打开LLD工具，启用宏"),
    ("2", "基础配置：产品类型 GaussDB_DWS，安装模式「集群扩容」，认证模式与原集群一致"),
    ("3", "选择用户：root或omm（与原集群一致）"),
    ("4", "输入扩容节点数量，自定义套餐可选"),
    ("5", "IP规划与进程部署：配置机架/管理IP/业务IP"),
    ("6", "节点信息：CPU核数/内存/磁盘数/主机名"),
    ("7", "磁盘配置：配置扩容节点数据盘"),
    ("8", "实例参数配置：确保集群CN总数 <= 10"),
    ("9", "生成配置文件 -> expandTemplet.xml, preinstall.ini, precheck配置"),
]
for i, (num, desc) in enumerate(lld_steps):
    y = 3.9 + i * 0.35
    add_rect(sl, 0.5, y, 0.35, 0.28, BLUE_DARK)
    add_tb(sl, 0.5, y + 0.02, 0.35, 0.25, num, sz=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tb(sl, 1.0, y, 11, 0.3, desc, sz=10, color=DARK)

# S7: Preinstall and Node Add
sl = S()
add_tb(sl, 0.5, 0.3, 12, 0.6, "3.2 Preinstall 与 节点添加", sz=26, bold=True, color=BLUE_DARK)
add_tb(sl, 0.5, 1.0, 12, 0.4, "◎新节点 Preinstall 步骤", sz=16, bold=True, color=BLUE)
add_ml(sl, 0.5, 1.5, 12, 2.5, [
    {"t": "① 拷贝 preset 目录到所有扩容节点 -> /opt/preset/preset.sh", "s": 11},
    {"t": "② 上传 software 目录到主OMS节点 /opt/FusionInsight_SetupTool/", "s": 11},
    {"t": "③ 确认镜像文件已挂载，数据盘已分区（g_parted=0 表示已分区）", "s": 11},
    {"t": "④ 配置 preinstall.ini -> g_hosts / g_user_name / g_port / g_parted_conf / g_pkgs_dir", "s": 11},
    {"t": "⑤ 可选：开启core dump -> g_core_dump=1, g_core_dump_dir=/var/log/core", "s": 11},
    {"t": "⑥ 执行 ./preinstall.sh -> 自动触发 precheck -> 确认无 ERROR", "s": 11},
    {"t": "⑦ 扩容前巡检：SysChecker 部署 -> 数据收集 -> 倾斜检查 -> 网速检查", "s": 11},
])
add_tb(sl, 0.5, 4.2, 12, 0.4, "◎模板添加节点（推荐方式）", sz=16, bold=True, color=BLUE)
add_ml(sl, 0.5, 4.7, 12, 2.5, [
    {"t": "① Manager -> [集群] -> [主机] -> [添加] -> [模板添加主机]", "s": 11},
    {"t": "② 选择 LLD生成的 扩容模板 -> 配置 root/omm 密码 -> [查找]", "s": 11},
    {"t": "③ 跳转至[确认]页 -> [上一步] -> [分配角色] -> 取消MS角色分配", "s": 11},
    {"t": "④ [提交] -> 等待扩容进度完成 -> 验证[运行状态]为[良好]", "s": 11},
    {"t": "⑤ Manager -> MPPDB -> [实例] -> [添加实例] -> 勾选新增实例", "s": 11},
    {"t": "⑥ 配置参数：Coordinator个数<=10，DN数据盘数与原节点一致", "s": 11},
    {"t": "⑦ [提交] -> 重启配置过期服务 -> 验证告警", "s": 11},
])

# S8: Data Redistribution
sl = S()
add_tb(sl, 0.5, 0.3, 12, 0.6, "4. 数据重分布策略", sz=26, bold=True, color=BLUE_DARK)
add_tb(sl, 0.5, 1.0, 12, 0.4, "◎重分布模式对比", sz=16, bold=True, color=BLUE)
add_tbl(sl, 0.5, 1.5, 12.3, 3.0, [
    ["对比项", "离线重分布", "在线重分布"],
    ["业务影响", "仅支持简单DQL，不支持DDL/DCL", "支持INSERT/UPDATE/DELETE及部分DDL"],
    ["锁等待超时", "默认20分钟写锁超时", "阻塞用户数据更新操作"],
    ["适用场景", "业务低峰期，可停机窗口", "需持续提供业务服务"],
    ["执行方式", "Manager一键触发", "Manager一键触发(mppdb.insert.expand=true)"],
    ["失败恢复", "可暂停/重试/调整并发", "可暂停/重试/调整并发"],
    ["推荐策略", "分段执行，控制并发数", "配合业务窗口，分批重分布"],
], col_w=[2, 3.5, 3.5])

add_tb(sl, 0.5, 4.8, 12, 0.4, "◎重分布最佳实践", sz=16, bold=True, color=BLUE)
add_ml(sl, 0.5, 5.3, 12, 2.0, [
    {"t": "• 关闭自动重分布 -> 手动触发，确保扩容和重分布均可失败重试", "s": 11},
    {"t": "• 支持技 Schema/表维度设置重分布顺序，暂停状态下可调整并发数", "s": 11},
    {"t": "• 重分布前执行 ANALYZE VERBOSE，保证统计信息准确性", "s": 11},
    {"t": "• 避免在重分布期间执行长时间查询（>20分钟），防止锁等待超时", "s": 11},
    {"t": "• 监控重分布进度：Manager -> MPPDB -> 更多 -> 重分布历史", "s": 11},
])

# S9: Risk Control
sl = S()
add_tb(sl, 0.5, 0.3, 12, 0.6, "5. 风险控制与应急预案", sz=26, bold=True, color=BLUE_DARK)
add_tb(sl, 0.5, 1.0, 12, 0.4, "◎扩容前必检项", sz=16, bold=True, color=BLUE)
add_tbl(sl, 0.5, 1.5, 12.3, 3.8, [
    ["检查项", "风险描述", "控制措施", "严重程度"],
    ["数据倾斜", "倾斜率过高导致重分布后各节点负载不均", "执行 CheckTableSkew，确保倾斜率<=10%", "高"],
    ["剩余容量", "重分布期间磁盘空间不足", "剩余容量 >= 最大单表1.5倍，不足则增加配额", "高"],
    ["网络带宽", "节点间网络瓶颈影响重分布效率", "执行 CheckNetSpeed，万兆网络互联", "中"],
    ["KV/系统表", "系统表脏页导致性能下降", "VACUUM FULL FREEZE PG_* 系列系统表", "中"],
    ["大表脏页率", "脏页率过高影响重分布", "倾斜率建议不超过20%~30%", "中"],
    ["集群状态", "集群不在balanced状态", "检查集群状态，确保balanced后再扩容", "高"],
    ["OS编码格式", "新节点locale不一致", "确保 en_US.UTF-8 一致", "低"],
    ["磁盘挂载", "新节点数据盘未分区", "确认分区且无数据，g_parted=0", "高"],
], col_w=[2.5, 3.5, 4.5, 1.3])

add_tb(sl, 0.5, 5.5, 12, 0.4, "◎故障处理矩阵", sz=16, bold=True, color=BLUE)
add_tbl(sl, 0.5, 6.0, 12.3, 1.3, [
    ["故障场景", "处理方式"],
    ["添加实例MPPDB报错（配置阶段）", "查看配置服务详情，确认配置与扩容规划一致后重试"],
    ["添加实例MPPDB报错（初始化阶段）", "查看 postinstall.log 日志定位原因，修复后重试"],
    ["任务界面误关闭", "Manager -> MPPDB -> 实例 -> 检查新实例信息，重试或删除主机后重新添加"],
    ["重分布过程中断", "支持暂停/重试/调整并发数，避免kill -9"],
    ["扩容期间告警", "参照故障处理步骤清除告警，必要时联系技术支持"],
], col_w=[4, 8])

# S10: Rollback Strategy
sl = S()
add_tb(sl, 0.5, 0.3, 12, 0.6, "6. 回滚方案", sz=26, bold=True, color=BLUE_DARK)
add_tb(sl, 0.5, 1.0, 12, 0.4, "◎各阶段回退策略", sz=16, bold=True, color=BLUE)
add_tbl(sl, 0.5, 1.5, 12.3, 2.8, [
    ["阶段", "故障场景", "回退操作"],
    ["扩容前", "环境检查不通过", "修复检查项后重试，不影响生产"],
    ["扩容-主机添加", "节点添加失败", "删除失败主机 -> 修复 -> 重新添加"],
    ["扩容-实例添加", "实例部署失败", "查看任务详情 -> 修复 -> 重试；或删除主机后重新添加"],
    ["重分布中", "重分布报错/中断", "支持暂停/重试/调整并发；不支持回滚到扩容前"],
    ["重分布后", "业务性能不达标", "评估是否需要回退（需重装集群，操作复杂）"],
    ["升级阶段", "升级失败/业务异常", "UpdateService -> 回滚 -> 回滚Preinstall -> 卸载UpdateService"],
], col_w=[2, 3, 7])

add_tb(sl, 0.5, 4.5, 12, 0.4, "◎关键回滚命令速查", sz=16, bold=True, color=BLUE)
add_ml(sl, 0.5, 5.0, 12, 2.0, [
    {"t": "• 回滚集群版本：UpdateService页面 -> [回滚] -> 重试失败操作或忽略", "s": 11},
    {"t": "• 回滚MPPDB文件：sh rollUpgradeMPP.sh -a <源版本> <目标版本>", "s": 11},
    {"t": "• 回滚Preinstall：cd /opt/FusionInsight_SetupTool/upgrade -> ./upgrade_setup.sh -p -r -n", "s": 11},
    {"t": "• 卸载UpdateService：cd /tmp/FusionInsight_UpdateService -> ./uninstall.sh", "s": 11},
    {"t": "• 解除容灾关系：python3 SyncDataToStby.py -t set-independent (源/目标集群各自执行)", "s": 11},
])

# S11: Key Risk Points
sl = S()
add_tb(sl, 0.5, 0.3, 12, 0.6, "5.2 扩容风险雷达图（关键风险TOP 6）", sz=24, bold=True, color=BLUE_DARK)
risks = [
    ("数据倾斜", "高", "扩容前必须通过CheckTableSkew检查，确保倾斜率<=10%，否则重分布后负载不均"),
    ("磁盘空间", "高", "剩余容量需>=最大单表的1.5倍，重分布期间临时空间需求剧增"),
    ("网络瓶颈", "中高", "万兆网络必须稳定，CheckNetSpeed通过后才可执行扩容"),
    ("锁冲突", "中", "重分布期间避免长查询（>20分钟），写锁默认超时20分钟"),
    ("版本兼容", "中", "跨MPPDB版本不支持容灾，扩容必须同版本同构"),
    ("时间窗口", "中", "重分布耗时与数据量正相关，需预留充足业务窗口"),
]
for i, (risk, level, desc) in enumerate(risks):
    y = 1.1 + i * 0.95
    lc = RED if level in ("高", "中高") else ORANGE
    add_rect(sl, 0.5, y, 1.5, 0.7, lc)
    add_tb(sl, 0.5, y + 0.1, 1.5, 0.35, risk, sz=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_tb(sl, 0.5, y + 0.4, 1.5, 0.3, "风险等级: " + level, sz=9, color=WHITE, align=PP_ALIGN.CENTER)
    add_tb(sl, 2.3, y + 0.1, 10, 0.6, desc, sz=11, color=DARK)

# S12: Project Timeline
sl = S()
add_tb(sl, 0.5, 0.3, 12, 0.6, "7. 项目实施计划", sz=26, bold=True, color=BLUE_DARK)
add_tb(sl, 0.5, 1.0, 12, 0.4, "◎阶段划分与时间线", sz=16, bold=True, color=BLUE)
add_tbl(sl, 0.5, 1.5, 12.3, 3.2, [
    ["阶段", "主要工作", "预计工期", "交付物"],
    ["T0: 环境准备", "软硬件信息收集、网络策略开通、LLD工具准备", "3-5天", "信息收集表、网络策略文档"],
    ["T1: 测试环境演练", "测试集群搭建、迁移演练、扩容演练、耗时评估", "5-7天", "演练报告、耗时记录"],
    ["T2: 生产扩容实施", "Preinstall -> 添加节点 -> 添加实例 -> 重分布", "2-3天", "扩容实施记录"],
    ["T3: 业务验证", "数据校验、业务功能测试、性能测试", "3-5天", "测试报告、验证签字"],
    ["T4: 割接上线", "DNS切换、客户端更新、监控上线", "1天", "割接报告、上线确认"],
    ["T5: 持续观察", "7天业务观察、性能监控、问题处理", "7天", "观察期报告、优化建议"],
], col_w=[2, 4, 2, 4.3])

add_tb(sl, 0.5, 5.0, 12, 0.4, "◎关键里程碑", sz=16, bold=True, color=BLUE)
milestones = [
    "M1: 测试环境搭建完成，功能验证通过",
    "M2: 生产环境扩容完成，数据重分布100%",
    "M3: 业务验证通过，客户签字确认",
    "M4: 正式割接上线，系统运行稳定",
]
for i, ms in enumerate(milestones):
    y = 5.5 + i * 0.4
    add_rect(sl, 0.5, y, 0.3, 0.3, GREEN)
    add_tb(sl, 1.0, y, 11, 0.3, ms, sz=12, color=DARK)

# S13: Summary
sl = S()
bg = sl.background.fill; bg.solid(); bg.fore_color.rgb = BLUE_DARK
add_tb(sl, 1, 1.5, 11, 0.8, "总结与建议", sz=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_ml(sl, 1.5, 2.8, 10, 3.5, [
    {"t": "1. 采用Roach双集群架构实现平滑迁移，全量+增量同步确保数据零丢失", "s": 14, "c": WHITE},
    {"t": "2. 扩容三阶段（迁移->升级->扩容）解耦执行，各阶段可独立回退", "s": 14, "c": WHITE},
    {"t": "3. 在线重分布模式支持业务不中断扩容，配合分段执行降低风险", "s": 14, "c": WHITE},
    {"t": "4. 完善的前置检查体系（数据倾斜/磁盘空间/网络）是扩容成功的基石", "s": 14, "c": WHITE},
    {"t": "5. 标准化回滚流程覆盖各故障场景，确保极端情况下的业务连续性", "s": 14, "c": WHITE},
    {"t": "6. 建议先进行完整测试环境演练，精确评估耗时后再执行生产操作", "s": 14, "c": WHITE},
])
add_tb(sl, 1, 5.5, 11, 0.5, "Thanks & Questions", sz=24, bold=True,
       color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.CENTER)

# Save
output_path = os.path.expanduser("~/Desktop/02-POC与售前/DWS-POC/稠州银行DWS扩容demo/稠州银行DWS扩容Demo.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
